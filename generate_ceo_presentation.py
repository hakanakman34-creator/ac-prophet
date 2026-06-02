import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def generate_ceo_pptx():
    output_path = r"c:\AC PROPHET\AC_PROPHET_CEO_Presentation.pptx"
    bg_img = r"C:\Users\TRON PCH\.gemini\antigravity\brain\2af92f41-e552-4ef6-96b7-8756516a1cae\ceo_presentation_bg_1780386169028.png"
    workflow_img = r"C:\Users\TRON PCH\.gemini\antigravity\brain\2af92f41-e552-4ef6-96b7-8756516a1cae\ai_workflow_diagram_1780386158350.png"
    hvac_img = r"C:\Users\TRON PCH\.gemini\antigravity\brain\2af92f41-e552-4ef6-96b7-8756516a1cae\ai_hvac_optimization_1780386145211.png"
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    WHITE = RGBColor(255, 255, 255)
    TEAL = RGBColor(0, 200, 255)
    NAVY = RGBColor(10, 25, 47)
    LIGHT_GREY = RGBColor(230, 235, 241)
    
    layout_blank = prs.slide_layouts[6]
    
    def apply_bg(slide):
        if os.path.exists(bg_img):
            slide.shapes.add_picture(bg_img, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        else:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = NAVY
            
    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(layout_blank)
    apply_bg(slide1)
    
    # Title overlay box
    overlay = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(5, 10, 20)
    # Using transparency via XML is tricky, we'll just use a solid dark color with no line
    overlay.line.fill.background()
    
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(11), Inches(2.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "AC PROPHET"
    p.font.name = "Segoe UI Light"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Driven Peak Season Operations & Dispatch Control Center"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(28)
    p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    
    # --- Slide 2: The Challenge & Vision ---
    slide2 = prs.slides.add_slide(layout_blank)
    apply_bg(slide2)
    
    if os.path.exists(hvac_img):
        slide2.shapes.add_picture(hvac_img, Inches(7.5), Inches(1.5), width=Inches(5), height=Inches(5))
        
    tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(6.5), Inches(5.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Strategic Vision & Value Proposition"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    points = [
        ("The Peak Season Challenge", "During extreme summer waves, manual dispatching fails to optimize workforce capacity, leading to accumulating backlogs and reduced customer satisfaction."),
        ("The AI Revolution", "AC PROPHET shifts operations from reactive firefighting to predictive, automated capacity management."),
        ("Operational Efficiency Reimagined", "Achieve up to 40% reduction in average wait times through intelligent load balancing and zero-touch tactical dispatching.")
    ]
    
    for title, desc in points:
        p_t = tf2.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEAL
        p_t.space_before = Pt(24)
        
        p_d = tf2.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(16)
        p_d.font.color.rgb = LIGHT_GREY
        p_d.space_before = Pt(6)

    # --- Slide 3: Multi-Agent Workflow ---
    slide3 = prs.slides.add_slide(layout_blank)
    apply_bg(slide3)
    
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "Next-Generation Multi-Agent Architecture"
    p3.font.size = Pt(36)
    p3.font.bold = True
    p3.font.color.rgb = WHITE
    
    if os.path.exists(workflow_img):
        # Place it nicely
        slide3.shapes.add_picture(workflow_img, Inches(0.8), Inches(1.5), width=Inches(6.5), height=Inches(5.5))
        
    tb_w = slide3.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(5), Inches(5))
    tf_w = tb_w.text_frame
    tf_w.word_wrap = True
    
    agents = [
        ("1. Forecaster Agent", "Fuses 7-day weather predictions with historical backlog trends to predict real incoming daily volumes."),
        ("2. Watchdog Agent", "Cross-references predictions against daily service capacities, identifying critical failure zones and wait-time bottlenecks."),
        ("3. Commander Agent", "Autonomously drafts tactical orders to shift resources from underutilized districts to critical zones before delays occur.")
    ]
    
    for ag, desc in agents:
        p_a = tf_w.add_paragraph()
        p_a.text = ag
        p_a.font.size = Pt(22)
        p_a.font.bold = True
        p_a.font.color.rgb = TEAL
        p_a.space_before = Pt(20)
        
        p_ad = tf_w.add_paragraph()
        p_ad.text = desc
        p_ad.font.size = Pt(14)
        p_ad.font.color.rgb = LIGHT_GREY
        p_ad.space_before = Pt(6)

    # --- Slide 4: Real-world Impact & Cloud Readiness ---
    slide4 = prs.slides.add_slide(layout_blank)
    apply_bg(slide4)
    
    tb4 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
    p4 = tb4.text_frame.paragraphs[0]
    p4.text = "Future-Proof & Cloud Ready"
    p4.font.size = Pt(36)
    p4.font.bold = True
    p4.font.color.rgb = WHITE
    p4.alignment = PP_ALIGN.CENTER
    
    impact_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.2), Inches(5.3), Inches(4))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = RGBColor(15, 30, 50)
    impact_box.line.color.rgb = TEAL
    
    ib_tf = impact_box.text_frame
    ib_tf.word_wrap = True
    p_i = ib_tf.paragraphs[0]
    p_i.text = "Operational Impact"
    p_i.font.size = Pt(24)
    p_i.font.color.rgb = TEAL
    p_i.alignment = PP_ALIGN.CENTER
    
    impacts = ["Eliminates manual guessing in dispatch", "Instantly visualizes crisis zones", "Provides actionable, written orders daily", "Maintains high CSAT during extreme heat"]
    for imp in impacts:
        pp = ib_tf.add_paragraph()
        pp.text = "• " + imp
        pp.font.size = Pt(16)
        pp.font.color.rgb = LIGHT_GREY
        pp.space_before = Pt(14)

    cloud_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.2), Inches(5.3), Inches(4))
    cloud_box.fill.solid()
    cloud_box.fill.fore_color.rgb = RGBColor(15, 30, 50)
    cloud_box.line.color.rgb = TEAL
    
    cb_tf = cloud_box.text_frame
    cb_tf.word_wrap = True
    p_c = cb_tf.paragraphs[0]
    p_c.text = "Google Cloud Roadmap"
    p_c.font.size = Pt(24)
    p_c.font.color.rgb = TEAL
    p_c.alignment = PP_ALIGN.CENTER
    
    clouds = ["Serverless execution with Cloud Run", "Cost-effective, scale-to-zero infrastructure", "Live Google Sheets API integrations", "100% Data persistence & enterprise security"]
    for cd in clouds:
        pp = cb_tf.add_paragraph()
        pp.text = "• " + cd
        pp.font.size = Pt(16)
        pp.font.color.rgb = LIGHT_GREY
        pp.space_before = Pt(14)

    # --- Slide 5: Thank You ---
    slide5 = prs.slides.add_slide(layout_blank)
    apply_bg(slide5)
    
    overlay2 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.5), Inches(9.333), Inches(2.5))
    overlay2.fill.solid()
    overlay2.fill.fore_color.rgb = RGBColor(5, 10, 20)
    overlay2.line.fill.background()
    
    tb5 = slide5.shapes.add_textbox(Inches(2), Inches(2.6), Inches(9.333), Inches(2.3))
    tf5 = tb5.text_frame
    p = tf5.paragraphs[0]
    p.text = "Transforming Operations, Today."
    p.font.name = "Segoe UI Light"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf5.add_paragraph()
    p2.text = "Thank You | AI Multi-Agent Group"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEAL
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    prs.save(output_path)
    print(f"Breathtaking CEO PPTX saved to {output_path}")

if __name__ == '__main__':
    generate_ceo_pptx()
