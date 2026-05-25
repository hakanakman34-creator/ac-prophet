import os
import json
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# ----------------- FONT REGISTRATION -----------------
SEG_REG = "C:\\Windows\\Fonts\\segoeui.ttf"
SEG_BOLD = "C:\\Windows\\Fonts\\segoeuib.ttf"
SEG_ITALIC = "C:\\Windows\\Fonts\\segoeuii.ttf"

pdfmetrics.registerFont(TTFont('SegoeUI', SEG_REG))
pdfmetrics.registerFont(TTFont('SegoeUIBold', SEG_BOLD))
pdfmetrics.registerFont(TTFont('SegoeUIItalic', SEG_ITALIC))

KO_REG = "C:\\Windows\\Fonts\\malgun.ttf"
KO_BOLD = "C:\\Windows\\Fonts\\malgunbd.ttf"

pdfmetrics.registerFont(TTFont('MalgunGothic', KO_REG))
pdfmetrics.registerFont(TTFont('MalgunGothicBold', KO_BOLD))

# ----------------- NUMBERED CANVAS -----------------
class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_decorations(num_pages)
            super().showPage()
        super().save()

    def draw_decorations(self, page_count):
        if self._pageNumber == 1:
            return
        
        self.saveState()
        is_ko = "KO" in self._filename.upper()
        font_name = "MalgunGothic" if is_ko else "SegoeUI"
        
        self.setFont(font_name, 9)
        self.setFillColor(colors.HexColor("#666666"))
        
        header_text = "AC PROPHET - Samsung HVAC Operations Control Center"
        self.drawString(54, 790, header_text)
        
        self.setStrokeColor(colors.HexColor("#CCCCCC"))
        self.setLineWidth(0.5)
        self.line(54, 782, 541, 782)
        
        self.line(54, 60, 541, 60)
        self.drawString(54, 45, "CONFIDENTIAL - Samsung HVAC Internal Use Only")
        page_lbl = f"{self._pageNumber} / {page_count}"
        self.drawRightString(541, 45, page_lbl)
        
        self.restoreState()


# ----------------- DYNAMIC MARKDOWN PARSER -----------------
def format_markdown_bold(text):
    # Format inline bold **text** to HTML <b>text</b>
    return re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text)

def render_table_flowable(table_lines, font_reg, font_bold, body_style):
    data = []
    for line in table_lines:
        cells = [c.strip() for c in line.split("|")[1:-1]]
        if not cells or all(c.startswith(":") or c.startswith("-") for c in cells):
            continue
        row_data = []
        for cell in cells:
            cell_html = format_markdown_bold(cell)
            row_data.append(Paragraph(cell_html, body_style))
        data.append(row_data)
        
    if not data:
        return Spacer(1, 1)
        
    num_cols = len(data[0])
    col_width = 487 / num_cols
    # Custom elegant widths for our 4-column system table
    col_widths = [110, 60, 180, 137] if num_cols == 4 else [col_width] * num_cols
    
    t = Table(data, colWidths=col_widths)
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#002B66')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#D3D3D3')),
    ]))
    
    # Update text color of header row in Table elements
    # Since headers are rendered as Paragraphs using body_style (dark color), we can override them if needed,
    # but the background #002B66 looks highly contrastive and extremely readable in Segoe UI.
    return t

def render_blockquote_flowable(quote_lines, font_reg, alert_style):
    full_text = " ".join(quote_lines)
    full_text = full_text.replace("[!WARNING]", "<b>⚠️ WARNING:</b>")
    full_text = full_text.replace("[!NOTE]", "<b>ℹ️ NOTE:</b>")
    full_text = full_text.replace("[!IMPORTANT]", "<b>❗ IMPORTANT:</b>")
    
    para = Paragraph(full_text, alert_style)
    alert_table = Table([[para]], colWidths=[487])
    
    is_warning = "WARNING" in full_text or "IMPORTANT" in full_text
    bg_color = colors.HexColor('#FFF2E6') if is_warning else colors.HexColor('#F8F9FA')
    border_color = colors.HexColor('#FF8000') if is_warning else colors.HexColor('#CCCCCC')
    
    alert_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), bg_color),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('LEFTPADDING', (0,0), (-1,-1), 12),
        ('RIGHTPADDING', (0,0), (-1,-1), 12),
        ('BOX', (0,0), (-1,-1), 1, border_color),
    ]))
    return alert_table

def generate_pdf_from_markdown(md_path, pdf_path, lang_code):
    is_ko = lang_code == "KO"
    font_reg = "MalgunGothic" if is_ko else "SegoeUI"
    font_bold = "MalgunGothicBold" if is_ko else "SegoeUIBold"
    
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=A4,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=72
    )
    
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'CoverTitle', parent=styles['Normal'], fontName=font_bold, fontSize=24, leading=30,
        textColor=colors.HexColor('#002B66'), spaceAfter=15
    )
    subtitle_style = ParagraphStyle(
        'CoverSubtitle', parent=styles['Normal'], fontName=font_reg, fontSize=12, leading=16,
        textColor=colors.HexColor('#1E51A4'), spaceAfter=25
    )
    meta_style = ParagraphStyle(
        'CoverMeta', parent=styles['Normal'], fontName=font_reg, fontSize=10, leading=14,
        textColor=colors.HexColor('#666666'), spaceAfter=80
    )
    h1_style = ParagraphStyle(
        'H1', parent=styles['Heading1'], fontName=font_bold, fontSize=16, leading=20,
        textColor=colors.HexColor('#002B66'), spaceBefore=18, spaceAfter=10, keepWithNext=True
    )
    h2_style = ParagraphStyle(
        'H2', parent=styles['Heading2'], fontName=font_bold, fontSize=13, leading=17,
        textColor=colors.HexColor('#1E51A4'), spaceBefore=14, spaceAfter=8, keepWithNext=True
    )
    body_style = ParagraphStyle(
        'Body', parent=styles['Normal'], fontName=font_reg, fontSize=10, leading=14,
        textColor=colors.HexColor('#2D3748'), spaceAfter=8
    )
    bullet_style = ParagraphStyle(
        'Bullet', parent=styles['Normal'], fontName=font_reg, fontSize=10, leading=14,
        textColor=colors.HexColor('#2D3748'), leftIndent=15, firstLineIndent=-10, spaceAfter=6
    )
    alert_style = ParagraphStyle(
        'Alert', parent=styles['Normal'], fontName=font_reg, fontSize=9, leading=13,
        textColor=colors.HexColor('#6A3B00'), spaceAfter=8
    )
    
    story = []
    
    # Metadata for cover
    title_text = ""
    subtitle_text = ""
    
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    # Read cover elements from first few lines
    for line in lines[:5]:
        stripped = line.strip()
        if stripped.startswith("# "):
            title_text = stripped[2:].replace("❄️ ", "").strip()
        elif stripped.startswith("## "):
            subtitle_text = stripped[3:].strip()
            
    # Cover Page
    story.append(Spacer(1, 80))
    story.append(Paragraph(title_text, title_style))
    story.append(Paragraph(subtitle_text, subtitle_style))
    
    line_table = Table([[""]], colWidths=[487], rowHeights=[4])
    line_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#002B66')),
        ('TOPPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
    ]))
    story.append(line_table)
    story.append(Spacer(1, 15))
    
    meta_desc = "AC PROPHET decision support matrix guide." if lang_code=="EN" else ("삼성 HVAC 의사결정 지원 매트릭스 매뉴얼" if is_ko else "Samsung HVAC karar destek matrisi kitapçığı.")
    story.append(Paragraph(meta_desc, meta_style))
    story.append(Spacer(1, 120))
    
    author_text = "AI Multi-Agent Group & Developer Team" if lang_code=="EN" else ("AI 멀티 에이전트 그룹 & 개발 엔지니어링 팀" if is_ko else "Yapay Zeka Çoklu Ajan Grubu & Geliştirici Ekibi")
    author_table_data = [
        [Paragraph(f"<b>Author:</b> {author_text}", body_style)],
        [Paragraph(f"<b>System Phase:</b> Operational Control / AI Strategic Model", body_style)],
        [Paragraph(f"<b>Documentation Level:</b> Comprehensive Detailed Manual", body_style)],
        [Paragraph(f"<b>Date:</b> May 2026", body_style)]
    ]
    author_table = Table(author_table_data, colWidths=[400])
    author_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#F8F9FA')),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('LEFTPADDING', (0,0), (-1,-1), 12),
        ('RIGHTPADDING', (0,0), (-1,-1), 12),
        ('BOX', (0,0), (-1,-1), 0.5, colors.HexColor('#E2E8F0')),
    ]))
    story.append(author_table)
    story.append(PageBreak())
    
    # Parse Markdown Contents Dynamically
    in_table = False
    table_lines = []
    in_quote = False
    quote_lines = []
    in_code = False
    
    for line in lines:
        stripped = line.strip()
        
        # Skip title and subtitle from the body to prevent duplicates
        if stripped.startswith("# ") or stripped.startswith("## ") and stripped == subtitle_text:
            continue
            
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue
            
        # Table parsing
        if stripped.startswith("|"):
            in_table = True
            table_lines.append(line)
            continue
        elif in_table:
            story.append(render_table_flowable(table_lines, font_reg, font_bold, body_style))
            story.append(Spacer(1, 10))
            table_lines = []
            in_table = False
            
        # Blockquote parsing
        if stripped.startswith(">"):
            in_quote = True
            quote_lines.append(stripped[1:].strip())
            continue
        elif in_quote:
            story.append(render_blockquote_flowable(quote_lines, font_reg, alert_style))
            story.append(Spacer(1, 10))
            quote_lines = []
            in_quote = False
            
        if not stripped:
            continue
            
        # Headers
        if stripped.startswith("## "):
            h_text = stripped[3:].strip()
            story.append(Paragraph(h_text, h1_style))
            story.append(Spacer(1, 8))
        elif stripped.startswith("### "):
            h_text = stripped[4:].strip()
            story.append(Paragraph(h_text, h2_style))
            story.append(Spacer(1, 6))
        # Lists
        elif stripped.startswith("* ") or stripped.startswith("- "):
            bullet_text = "&bull; " + stripped[2:].strip()
            bullet_html = format_markdown_bold(bullet_text)
            story.append(Paragraph(bullet_html, bullet_style))
        elif stripped.startswith("1. ") or stripped.startswith("2. ") or stripped.startswith("3. ") or stripped.startswith("4. "):
            numbered_html = format_markdown_bold(stripped)
            story.append(Paragraph(numbered_html, bullet_style))
        # Separator line
        elif stripped == "---":
            story.append(Spacer(1, 10))
            line_table = Table([[""]], colWidths=[487], rowHeights=[0.5])
            line_table.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#CCCCCC')),
                ('TOPPADDING', (0,0), (-1,-1), 0),
                ('BOTTOMPADDING', (0,0), (-1,-1), 0),
            ]))
            story.append(line_table)
            story.append(Spacer(1, 10))
        # Plain text
        else:
            para_html = format_markdown_bold(stripped)
            story.append(Paragraph(para_html, body_style))
            
    # Trailing block flushes
    if table_lines:
        story.append(render_table_flowable(table_lines, font_reg, font_bold, body_style))
    if quote_lines:
        story.append(render_blockquote_flowable(quote_lines, font_reg, alert_style))
        
    doc.build(story, canvasmaker=NumberedCanvas)
    print(f"Compiled PDF directly from Markdown: {pdf_path}")


# ----------------- PowerPoint EXECUTIVE SLIDES -----------------
SLIDES_DATA = {
    "TR": {
        "title": "AC PROPHET",
        "subtitle": "Samsung HVAC Marmara Sevk ve Kapasite Optimizasyonu",
        "author": "Samsung HVAC İş Geliştirme & AI Ajan Grubu",
        "slides": [
            ("1. Giriş ve Vizyon (Executive Summary)", [
                ("Yüksek Sezon İhtiyaçları", [
                    "Peak Season döneminde aşırı yükselen klima servis taleplerini kontrol altına almak.",
                    "Sıcaklık ve hava durumu değişimleriyle servis taleplerini anlık korele etmek.",
                    "Yapay zeka ile manuel sevk süreçlerini tamamen otomatikleştirmek."
                ]),
                ("Çoklu Ajan Karar Desteği", [
                    "Hata payını en aza indiren 3 bağımsız uzman yapay zeka ajanı.",
                    "Eritilemeyen carryover (backlog) yükünün zincirleme bir sonraki güne devri."
                ])
            ]),
            ("2. Çoklu Ajan Mimarisi", [
                ("Forecaster Agent (Tahmin)", [
                    "Hava tahmin raporundaki sıcaklık artışları ve geçmiş iş yükü trendlerini birleştirir.",
                    "Her servis merkezi için 7 günlük yeni gelen ve sonraki güne sarkan iş tahminlerini üretir."
                ]),
                ("Watchdog & Commander", [
                    "Watchdog Agent tahmini iş yükünü servis kapasiteleriyle oranlayıp risk haritası çizer.",
                    "Commander Agent bekleme süresi 4 günü aşan kritik servisler için ekip sevk emirlerini hazırlar."
                ])
            ]),
            ("3. Operations Dashboard (Operasyon Paneli)", [
                ("Hava Durumu & Risk Haritası", [
                    "Marmara Bölgesi illerinin 7 günlük sıcaklık tahminleri listelenir.",
                    "Risk Haritası ilçeleri bekleme süresine göre yeşil, sarı ve kırmızı boyar.",
                    "Oynat/Duraklat butonu ile günden güne değişimi canlı animasyon olarak gösterir."
                ]),
                ("Commander Emirleri", [
                    "Ajanın o gün için ürettiği yazılı sevk ve ekip planı alt panelde görüntülenir."
                ])
            ]),
            ("4. Data Management (Veri Yönetimi)", [
                ("Manuel ve Toplu Giriş", [
                    "Seçilen tarih için servis bazlı iş verileri manuel tablo üzerinden girilir.",
                    "Excel şablonuna uygun toplu veri kütlesi sisteme tek seferde yüklenir."
                ]),
                ("Geçmiş Veri Düzenleyici (Data Fixer)", [
                    "Excel tablosunun tamamı arayüzde listelenir. Hatalı geçmiş hücreler elle düzenlenip kaydedilir."
                ])
            ]),
            ("5. Admin & Service District Mapping", [
                ("Kapasite Yönetimi", [
                    "Servis merkezlerinin aktif ekip sayısı ve ekip başına günlük iş tamamlama kapasitesi yönetilir."
                ]),
                ("İlçe Eşleşmeleri (Mapping)", [
                    "Hangi ilçenin hangi servis merkezine bağlı olduğu interaktif tablodan atanır.",
                    "Yapılan değişiklikler 'Kaydet ve Uygula' denildiğinde anında haritaya yansır."
                ])
            ]),
            ("6. Veri Modeli ve Dosya Düzeni", [
                ("Jobsdata.xlsx & ServiceList.xlsx", [
                    "Jobsdata.xlsx geçmiş iş yükü, tamamlanan işler ve backlog tarihçesini tutar.",
                    "ServiceList.xlsx servis tanımlı kapasite ve aktif ekip sayısı verilerini tutar."
                ]),
                ("İlçe Eşleşme Dosyaları", [
                    "marmara_ilce_listesi.xlsx ilçe atama kodlarını tutan ana tablodur.",
                    "service_district_map.json hızlı harita render'ı için Excel'den derlenen dosyadır."
                ])
            ]),
            ("7. Buluta Taşıma ve Google Cloud Planı", [
                ("Google Cloud Run & Docker", [
                    "Uygulamayı Dockerize ederek Cloud Run'da en düşük maliyetle barındırabiliriz.",
                    "Yalnızca kullanıldığında çalışan, SSL sertifikalı güvenli bulut yapısı."
                ]),
                ("Google Sheets API Entegrasyonu", [
                    "Excel dosyalarını Google E-Tablolar'a bağlayarak çoklu veri girişine açma.",
                    "Sunucu çökmesi veya sıfırlanmasında kalıcı veri koruması."
                ])
            ])
        ]
    },
    "EN": {
        "title": "AC PROPHET",
        "subtitle": "Samsung HVAC Marmara Dispatch & Capacity Optimization",
        "author": "Samsung HVAC Business Development & AI Agent Group",
        "slides": [
            ("1. Introduction and Vision (Executive Summary)", [
                ("Widescreen Dispatch Goals", [
                    "Proactive control over extremely high HVAC service demands during summer peak seasons.",
                    "Correlating temperature spikes with incoming call rates in real-time.",
                    "Replacing manual planning with automated multi-agent AI systems."
                ]),
                ("AI Strategic Advantage", [
                    "Three independent specialist agents negotiating capacity limits.",
                    "Accurate daily carryover accumulation models mapping real backlog flow."
                ])
            ]),
            ("2. Multi-Agent Architecture", [
                ("Forecaster Agent (Forecast)", [
                    "Merges 7-day weather predictions with historical request patterns.",
                    "Generates sifting carryover and new incoming workloads for each service center."
                ]),
                ("Watchdog & Commander", [
                    "Watchdog matches workload with daily limits and calculates district wait times.",
                    "Commander issues transfer orders when wait times exceed threshold (4 days)."
                ])
            ]),
            ("3. Operations Dashboard Section Guide", [
                ("Weather & Dynamic Maps", [
                    "Lists 7-day average temperatures for major Marmara cities.",
                    "Dynamic choropleth map plots district wait periods in Green, Yellow, and Red.",
                    "Interactive timeline playback animates backlog growth day by day."
                ]),
                ("Operational Action Plan", [
                    "Written commander actions are displayed clearly at the bottom of the screen."
                ])
            ]),
            ("4. Data Management Overview", [
                ("Manual & Bulk Pipelines", [
                    "Add new daily records containing assigned, completed and sifting jobs.",
                    "Seamlessly append bulk data files via uploader template frames."
                ]),
                ("Interactive Data Fixer", [
                    "Correct past historical logs directly using cell editors inside the browser tab."
                ])
            ]),
            ("5. Admin Tab & Service District Mapping", [
                ("Capacity Calibrator", [
                    "Fine-tune active Team Quantities and daily completion capacities per ASC."
                ]),
                ("District Matrix Mapping", [
                    "Assign districts to specific services directly within the screen grid.",
                    "Instantly regenerates the map visualizer configurations upon saving."
                ])
            ]),
            ("6. Data Model and File Register", [
                ("Jobsdata.xlsx & ServiceList.xlsx", [
                    "Jobsdata.xlsx tracks histories, backlog sizes and resolve rates.",
                    "ServiceList.xlsx maintains active teams and baseline completion capacities."
                ]),
                ("Lookup Configurations", [
                    "marmara_ilce_listesi.xlsx serves as the primary district mapping table.",
                    "service_district_map.json handles fast rendering mappings in the background."
                ])
            ]),
            ("7. Cloud Deployment & GCP Roadmap", [
                ("Google Cloud Run & Containers", [
                    "Dockerize the application and run serverless at minimal operational cost.",
                    "Scales down to zero on idle; secure HTTPS SSL configurations."
                ]),
                ("Live Google Sheets API", [
                    "Migrate spreadsheet files to Google Drive to enable secure multi-user editing.",
                    "Complete protection against server crashes or local storage loss."
                ])
            ])
        ]
    },
    "KO": {
        "title": "AC PROPHET",
        "subtitle": "삼성 HVAC 마르마라 서비스 배정 및 용량 최적화 시스템",
        "author": "삼성 HVAC 비즈니스 운영 및 AI 에이전트 그룹",
        "slides": [
            ("1. 도입 및 비전 (Executive Summary)", [
                ("성수기 운영 통제", [
                    "여름철 성수기 동안 극심하게 팽창하는 HVAC 수리 설치 수요를 선제 관리합니다.",
                    "기온 상승과 서비스 접수량의 상관관계를 실시간으로 자동 연산합니다.",
                    "지능형 에이전트 시스템을 적용하여 수동 배정 방식을 고도 자동화합니다."
                ]),
                ("멀티 에이전트 의사결정", [
                    "작업 분담 및 한도 조율을 주도하는 3개의 특화 에이전트로 구성됩니다.",
                    "처리 지연으로 남은 잔여 작업(Carryover)의 일차별 누적 모델을 완벽히 모사합니다."
                ])
            ]),
            ("2. 멀티 에이전트 아키텍처", [
                ("예측 에이전트 (Forecaster)", [
                    "7일간의 기온 예보 정보와 과거 누적 수주 트렌드를 융합합니다.",
                    "센터별 일차별 신규 배정 물량과 다음 날로 이월되는 대기 물량을 정밀 추론합니다."
                ]),
                ("감시 및 사령 에이전트", [
                    "감시 에이전트(Watchdog)는 대기 시간을 계산하고 구역별 3단계 경보를 발령합니다.",
                    "사령 에이전트(Commander)는 대기 4일 초과 위기 지점에 대해 지원 배치 명령을 생성합니다."
                ])
            ]),
            ("3. Operations Dashboard (운영 대시보드)", [
                ("기상 및 동적 리스크 맵", [
                    "마르마라 광역 도시별 7일 평균 기상 통계 그래프를 시각화합니다.",
                    "예상 지연율에 기초하여 각 군/구를 녹색, 황색, 적색으로 매핑합니다.",
                    "재생 애니메이션을 가동하여 7일간의 백로그 증감 과정을 한눈에 확인합니다."
                ]),
                ("사령 전술 보고서", [
                    "사령 에이전트가 발행한 지시사항이 하단 전술창에 실시간 바인딩됩니다."
                ])
            ]),
            ("4. 데이터 관리 (Data Management)", [
                ("수동 및 대량 통합 파이프라인", [
                    "신규 배정, 취소 및 처리 지표를 입력하면 대기 수량을 자동 파싱합니다.",
                    "업로더 프레임을 통해 대량의 엑셀 데이터셋을 고속 병합합니다."
                ]),
                ("인터랙티브 Data Fixer", [
                    "과거 적재된 전체 이력을 브라우저에서 즉시 타이핑 수정하여 반영합니다."
                ])
            ]),
            ("5. 관리자 탭 및 Service District Mapping", [
                ("서비스 용량 세팅", [
                    "서비스 지점별 가동 차량 수량과 일일 한도 처릿수를 정밀 제어합니다."
                ]),
                ("구역 매핑 (Service District Mapping)", [
                    "행정 구역에 매칭될 서비스 코드를 브라우저 화면에서 직관적으로 수정합니다.",
                    "저장 시 자동으로 백그라운드의 매핑 설정 엑셀을 저장하고 지도의 JSON을 재생성합니다."
                ])
            ]),
            ("6. 데이터 구조 및 시스템 파일 규칙", [
                ("Jobsdata.xlsx & ServiceList.xlsx", [
                    "Jobsdata.xlsx는 지점별 역사적 누적 이력, 대기 잔량 및 완료 지표들을 저장합니다.",
                    "ServiceList.xlsx는 모든 지점의 기동 팀 수, 기초 완료 한도를 담고 있습니다."
                ]),
                ("매핑 마스터 파일", [
                    "marmara_ilce_listesi.xlsx는 구역과 지점을 연계하는 마스터 데이터시트입니다.",
                    "service_district_map.json은 지도 시각화의 초고속 파싱을 지원하는 설정 파일입니다."
                ])
            ]),
            ("7. 클라우드 이관 및 Google Cloud 로드맵", [
                ("Google Cloud Run & Docker", [
                    "도커화된 앱을 Cloud Run 환경에 올려 서버리스 고속 호스팅을 개시합니다.",
                    "접속할 때만 자원을 소모하므로 운영비가 극적으로 절감되며 안전한 SSL이 탑재됩니다."
                ]),
                ("Google Sheets API 연동", [
                    "엑셀을 구글 실시간 E-Tablolar로 전환하여 실시간 협업 환경을 구축합니다.",
                    "서버가 초기화되어도 절대로 유실되지 않는 확실한 보안 persistence 제공."
                ])
            ])
        ]
    }
}

def generate_pptx(lang_code, output_path):
    lang_data = SLIDES_DATA[lang_code]
    font_name = "Segoe UI" if lang_code != "KO" else "Malgun Gothic"
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Executive Colors
    NAVY = RGBColor(15, 30, 54)
    SAMSUNG_BLUE = RGBColor(30, 81, 164)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GREY = RGBColor(244, 246, 248)
    DARK_GREY = RGBColor(74, 85, 104)
    BORDER_GREY = RGBColor(226, 232, 240)
    
    slide_layout = prs.slide_layouts[6]
    
    def format_tf(tf):
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        
    def add_bullet_run(tf, text, size=13, bold=False, color=DARK_GREY, space=6):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = "• " + text
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(space)
        p.line_spacing = 1.15

    # Title Slide
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = NAVY
    
    accent = slide1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = SAMSUNG_BLUE
    accent.line.fill.background()
    
    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lang_data["title"]
    p.font.name = font_name
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = lang_data["subtitle"]
    p2.font.name = font_name
    p2.font.size = Pt(20)
    p2.font.color.rgb = SAMSUNG_BLUE
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = f"Samsung HVAC Decision Support Dashboard | {lang_data['author']}"
    p3.font.name = font_name
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(160, 174, 192)
    p3.space_before = Pt(45)

    # Content Slides
    for slide_title, sections in lang_data["slides"]:
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GREY
        
        top_line = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = SAMSUNG_BLUE
        top_line.line.fill.background()
        
        header_tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        p_h = header_tb.text_frame.paragraphs[0]
        p_h.text = slide_title
        p_h.font.name = font_name
        p_h.font.size = Pt(28)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY
        
        num_cards = len(sections)
        if num_cards == 2:
            card_width = Inches(5.6)
            spacing = Inches(0.5)
            left_margin = Inches(0.8)
        elif num_cards == 3:
            card_width = Inches(3.6)
            spacing = Inches(0.46)
            left_margin = Inches(0.8)
        else:
            card_width = Inches(11.7)
            spacing = Inches(0)
            left_margin = Inches(0.8)
            
        for idx, (card_title, card_pts) in enumerate(sections):
            x_left = left_margin + idx * (card_width + spacing)
            
            card = slide.shapes.add_shape(1, x_left, Inches(1.4), card_width, Inches(5.3))
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = BORDER_GREY
            card.line.width = Pt(1.5)
            
            c_tb = slide.shapes.add_textbox(x_left + Inches(0.15), Inches(1.6), card_width - Inches(0.3), Inches(4.9))
            tf_c = c_tb.text_frame
            format_tf(tf_c)
            
            p_t = tf_c.paragraphs[0]
            p_t.text = card_title
            p_t.font.name = font_name
            p_t.font.size = Pt(18 if num_cards <= 2 else 15)
            p_t.font.bold = True
            p_t.font.color.rgb = SAMSUNG_BLUE
            p_t.space_after = Pt(12)
            
            for pt in card_pts:
                add_bullet_run(tf_c, pt, size=(13 if num_cards <= 2 else 11), space=(8 if num_cards <= 2 else 6))

    # Slide 8: Thank You
    slide8 = prs.slides.add_slide(slide_layout)
    slide8.background.fill.solid()
    slide8.background.fill.fore_color.rgb = NAVY
    
    accent8 = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent8.fill.solid()
    accent8.fill.fore_color.rgb = SAMSUNG_BLUE
    accent8.line.fill.background()
    
    tb8 = slide8.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(3.5))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    
    p_t8 = tf8.paragraphs[0]
    p_t8.text = "Thank You!" if lang_code != "KO" else "감사합니다!"
    p_t8.font.name = font_name
    p_t8.font.size = Pt(64)
    p_t8.font.bold = True
    p_t8.font.color.rgb = WHITE
    
    p_s8 = tf8.add_paragraph()
    p_s8.text = "AC PROPHET - Operational Stability & Capacity Dispatch Model" if lang_code != "KO" else "AC PROPHET - 안정적 운영 및 최적의 인력 자원 배치 모델"
    p_s8.font.name = font_name
    p_s8.font.size = Pt(18)
    p_s8.font.color.rgb = SAMSUNG_BLUE
    p_s8.space_before = Pt(16)
    
    prs.save(output_path)
    print(f"Generated PowerPoint: {output_path}")


# ----------------- MAIN EXECUTION PIPELINE -----------------
if __name__ == "__main__":
    out_dir = r"c:\AC PROPHET"
    os.makedirs(out_dir, exist_ok=True)
    
    # 1. Clean and Delete Old Files
    old_files = [
        "AC_PROPHET_User_Manual_TR.pdf", "AC_PROPHET_User_Manual_EN.pdf", "AC_PROPHET_User_Manual_KO.pdf",
        "AC_PROPHET_Presentation_TR.pptx", "AC_PROPHET_Presentation_EN.pptx", "AC_PROPHET_Presentation_KO.pptx"
    ]
    for f in old_files:
        path = os.path.join(out_dir, f)
        if os.path.exists(path):
            try:
                os.remove(path)
                print(f"Cleaned old file: {f}")
            except Exception as e:
                print(f"Error cleaning {f}: {e}")
                
    # 2. Compile PDFs directly from local markdown manual files!
    generate_pdf_from_markdown(
        os.path.join(out_dir, "ac_prophet_user_manual.md"),
        os.path.join(out_dir, "AC_PROPHET_User_Manual_TR.pdf"),
        "TR"
    )
    generate_pdf_from_markdown(
        os.path.join(out_dir, "ac_prophet_user_manual_en.md"),
        os.path.join(out_dir, "AC_PROPHET_User_Manual_EN.pdf"),
        "EN"
    )
    generate_pdf_from_markdown(
        os.path.join(out_dir, "ac_prophet_user_manual_ko.md"),
        os.path.join(out_dir, "AC_PROPHET_User_Manual_KO.pdf"),
        "KO"
    )
    
    # 3. Compile PPTX with custom premium card templates
    generate_pptx("TR", os.path.join(out_dir, "AC_PROPHET_Presentation_TR.pptx"))
    generate_pptx("EN", os.path.join(out_dir, "AC_PROPHET_Presentation_EN.pptx"))
    generate_pptx("KO", os.path.join(out_dir, "AC_PROPHET_Presentation_KO.pptx"))
    
    print("\nALL DYNAMIC MARKDOWN-TO-PDF FILES AND EXECUTIVE SLIDES RE-GENERATED SUCCESSFULLY!")
